#!/usr/bin/env python3
"""vISA 입문 PPT 렌더러 — python-pptx 기반.

슬라이드 스펙(JSON)을 받아 16:9 pptx 로 그린다.
내용량에 맞춰 글자 크기를 자동으로 줄이고, 한글은 East-Asian 타이프페이스로 지정한다.
"""
import math
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- 테마
INK = RGBColor(0x16, 0x20, 0x2C)
MUTED = RGBColor(0x5B, 0x6B, 0x7B)
ACCENT = RGBColor(0x0B, 0x5F, 0xC0)
ACCENT_D = RGBColor(0x08, 0x3F, 0x80)
AMBER = RGBColor(0xB4, 0x54, 0x09)
AMBER_BG = RGBColor(0xFF, 0xF7, 0xE6)
SOFT = RGBColor(0xF1, 0xF5, 0xF9)
LINE = RGBColor(0xD6, 0xDE, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x0F, 0x1B, 0x2A)
CODE_FG = RGBColor(0xE6, 0xED, 0xF3)
GREEN = RGBColor(0x1A, 0x7F, 0x37)
RED = RGBColor(0xB4, 0x22, 0x2A)

KO = "맑은 고딕"
LAT = "Segoe UI"
MONO = "Consolas"

OVERFLOW = []   # (쪽번호, 제목, 초과비율) — 최소 글자크기로도 안 들어간 슬라이드

SW, SH = Inches(13.333), Inches(7.5)
ML, MR = Inches(0.62), Inches(0.62)
BODY_W = SW - ML - MR
TITLE_Y = Inches(0.42)
BODY_Y = Inches(1.58)
BODY_BOT = Inches(6.86)
BODY_H = BODY_BOT - BODY_Y


def E(v):
    """EMU 는 정수여야 한다(ST_Coordinate = xsd:long).

    계산 중 나온 실수를 그대로 쓰면 x="2670596.6" 같은 값이 XML 에 박히고,
    PowerPoint 가 파일을 손상으로 판정한다. 모든 좌표는 이 함수를 통과시킨다.
    """
    return int(round(float(v)))


def vlen(s):
    """반각 단위 길이 (한글=2, 그 외=1)."""
    return sum(2 if ord(c) > 0x2E80 else 1 for c in s)


# ---------------------------------------------------------------- 저수준
def _ea(run, face):
    """East-Asian / complex-script 타이프페이스 지정.

    DrawingML CT_TextCharacterProperties 는 자식 순서를 latin -> ea -> cs 로 강제한다.
    순서를 어기면 PowerPoint 가 파일을 손상으로 보고 복구를 요구한다(LibreOffice 는 통과).
    """
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = etree.SubElement(rPr, qn("a:cs"))
    if latin is not None:
        latin.addnext(cs)   # -> latin, cs
        latin.addnext(ea)   # -> latin, ea, cs
    ea.set("typeface", face)
    cs.set("typeface", face)


def style(run, size=18, bold=False, color=INK, mono=False, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = MONO if mono else LAT
    _ea(run, MONO if mono else KO)


def code_size(lines, box_h_emu):
    """코드 글자 크기 — 가장 긴 줄이 폭에 들어가고 줄 수가 높이에 들어가는 최대 크기."""
    maxw = max((vlen(l) for l in lines), default=10)
    for cand in (15, 14, 13, 12, 11, 10, 9, 8, 7.5):
        if maxw * cand * 0.55 <= (BODY_W / 12700.0) - 26 and \
           len(lines) * cand * 1.32 <= (box_h_emu / 12700.0) - 16:
            return cand
    return 7.5


def code_rows(size, box_h_emu):
    """그 크기에서 코드 박스에 들어가는 최대 줄 수."""
    return max(1, int(((box_h_emu / 12700.0) - 16) / (size * 1.32)))


def tbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tf


def _runs(p, text, size, bold, color, mono):
    """줄바꿈은 <a:br/> 로 넣는다.

    run 의 text 에 '\n' 을 그대로 넣으면 <a:t> 안에 날 개행이 들어가는데,
    DrawingML 에서 그것은 줄바꿈이 아니라서 PowerPoint 가 한 줄로 이어 그린다.
    """
    for i, seg in enumerate(str(text).split("\n")):
        if i:
            p.add_line_break()
        r = p.add_run()
        r.text = seg
        style(r, size, bold, color, mono)


def para(tf, first, text, size, bold=False, color=INK, mono=False,
         space_before=0, space_after=0, indent=0, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = 1.18
    if indent:
        p.left_indent = Emu(indent)
    _runs(p, text, size, bold, color, mono)
    return p


def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=1.0):
    s = slide.shapes.add_shape(shape, E(x), E(y), E(w), E(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return s


def label(shape, text, size=12, bold=False, color=INK, mono=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = align
    _runs(p, text, size, bold, color, mono)


def arrow(slide, x1, y1, x2, y2, color=ACCENT, w=1.75):
    from pptx.enum.shapes import MSO_CONNECTOR
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1), E(y1), E(x2), E(y2))
    c.line.color.rgb = color
    c.line.width = Pt(w)
    ln = c.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return c


def hline(slide, x, y, w, color=LINE, lw=1.0):
    from pptx.enum.shapes import MSO_CONNECTOR
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x), E(y), E(x + w), E(y))
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    return c


# ---------------------------------------------------------------- 슬라이드 골격
class Deck:
    def __init__(self, title, subtitle, footer):
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self.blank = self.prs.slide_layouts[6]
        self.footer = footer
        self.n = 0
        self.part = ""
        self._cur_title = ""
        self._cover(title, subtitle)

    def _new(self):
        return self.prs.slides.add_slide(self.blank)

    def _chrome(self, s, title, subtitle=None, tight=False):
        """상단 제목 + 밑줄 + 하단 푸터/쪽번호.

        tight=True 는 도해 슬라이드용. 제목 영역을 0.3in 줄여 그림에 넘긴다 —
        도해는 폭이 넓을수록 안쪽 글자가 커지므로 세로 공간이 곧 가독성이다.
        """
        self._cur_title = title
        rect(s, 0, 0, Inches(0.16), SH, fill=ACCENT)
        if self.part:
            tf = tbox(s, ML, Inches(0.30 if not tight else 0.24), BODY_W, Inches(0.26))
            para(tf, True, self.part, 11.5, True, ACCENT)
        ty = (TITLE_Y if not tight else Inches(0.34)) + \
            (Inches(0.22 if not tight else 0.18) if self.part else 0)
        tsize = 30 if vlen(title) <= 46 else (26 if vlen(title) <= 64 else 22)
        if tight:
            tsize = min(tsize, 24)
        tf = tbox(s, ML, ty, BODY_W, Inches(0.78 if not tight else 0.52))
        para(tf, True, title, tsize, True, INK)
        sy = ty + Inches(0.72 if not tight else 0.50)
        if subtitle:
            tf = tbox(s, ML, sy, BODY_W, Inches(0.3))
            para(tf, True, subtitle, 14 if not tight else 12.5, False, MUTED)
            sy += Inches(0.34 if not tight else 0.28)
        hline(s, ML, sy + Inches(0.06), BODY_W, LINE, 1.25)
        self.n += 1
        tf = tbox(s, ML, SH - Inches(0.44), BODY_W - Inches(1.5), Inches(0.24))
        para(tf, True, self.footer, 9.5, False, MUTED)
        tf2 = tbox(s, SW - MR - Inches(1.2), SH - Inches(0.44), Inches(1.2), Inches(0.24))
        para(tf2, True, str(self.n), 9.5, False, MUTED, align=PP_ALIGN.RIGHT)
        return sy + Inches(0.30)

    # ------------------------------------------------------------ 표지
    def _cover(self, title, subtitle):
        s = self._new()
        rect(s, 0, 0, SW, SH, fill=ACCENT_D)
        rect(s, 0, SH - Inches(1.9), SW, Inches(1.9), fill=RGBColor(0x06, 0x33, 0x69))
        tf = tbox(s, Inches(1.0), Inches(2.1), SW - Inches(2.0), Inches(1.5))
        para(tf, True, title, 46, True, WHITE)
        tf = tbox(s, Inches(1.0), Inches(3.7), SW - Inches(2.0), Inches(1.0))
        for i, ln in enumerate(subtitle.split("\n")):
            para(tf, i == 0, ln, 19, False, RGBColor(0xC9, 0xDE, 0xF7), space_after=4)
        tf = tbox(s, Inches(1.0), SH - Inches(1.35), SW - Inches(2.0), Inches(0.8))
        para(tf, True, self.footer, 13, False, RGBColor(0x9E, 0xC1, 0xEB))

    # ------------------------------------------------------------ 부 표지
    def section(self, num, title, subtitle=""):
        self.part = ""
        s = self._new()
        rect(s, 0, 0, SW, SH, fill=SOFT)
        rect(s, 0, 0, Inches(0.16), SH, fill=ACCENT)
        nsz = 84 if vlen(num) <= 2 else (60 if vlen(num) <= 4 else 44)
        tf = tbox(s, Inches(1.1), Inches(2.35), Inches(2.6), Inches(1.4))
        para(tf, True, num, nsz, True, RGBColor(0xC3, 0xD4, 0xE8))
        tf = tbox(s, Inches(3.3), Inches(2.55), SW - Inches(4.4), Inches(1.0))
        para(tf, True, title, 36, True, INK)
        if subtitle:
            tf = tbox(s, Inches(3.3), Inches(3.62), SW - Inches(4.4), Inches(1.6))
            for i, ln in enumerate(subtitle.split("\n")):
                para(tf, i == 0, ln, 16, False, MUTED, space_after=5)
        return s

    # ------------------------------------------------------------ 큰 문장
    def big(self, text, sub=""):
        s = self._new()
        top = self._chrome(s, " ")
        tf = tbox(s, Inches(1.2), Inches(2.4), SW - Inches(2.4), Inches(2.2), MSO_ANCHOR.MIDDLE)
        size = 40 if vlen(text) <= 40 else (32 if vlen(text) <= 70 else 26)
        para(tf, True, text, size, True, ACCENT_D, align=PP_ALIGN.CENTER)
        if sub:
            tf = tbox(s, Inches(1.2), Inches(4.6), SW - Inches(2.4), Inches(1.0))
            para(tf, True, sub, 16, False, MUTED, align=PP_ALIGN.CENTER)
        return s

    # ------------------------------------------------------------ 자동 글자크기
    @staticmethod
    def _fit(items, width_emu, height_emu, sizes=(19, 18, 17, 16, 15, 14, 13, 12, 11, 10, 9)):
        w_pt = width_emu / 12700.0
        h_pt = height_emu / 12700.0
        for s in sizes:
            total = 0.0
            for it in items:
                lv = it.get("lv", 0) or 0
                avail = w_pt - lv * 0.30 * 72 - 14
                per = max(10.0, avail / (s * 0.5))
                lines = max(1, math.ceil(vlen(it["t"]) / per))
                fs = s if lv == 0 else s - 1.2
                total += lines * fs * 1.30 + (7 if lv == 0 else 3)
            if total <= h_pt:
                return s, total / h_pt
        return sizes[-1], total / h_pt

    # 도해 슬라이드의 글머리 기호 크기는 고정한다 (L-30).
    # _fit() 으로 슬라이드마다 다른 크기를 고르면 11/14/19pt 로 제각각이 된다 (실측 2026-08-14).
    DIA_ITEM_PT = 12

    @staticmethod
    def _needed_pt(items, width_emu, size):
        """이 글자 크기로 넣으면 세로로 몇 pt 가 필요한가."""
        w_pt = width_emu / 12700.0
        total = 0.0
        for it in items:
            lv = it.get("lv", 0) or 0
            avail = w_pt - lv * 0.30 * 72 - 14
            per = max(10.0, avail / (size * 0.5))
            lines = max(1, math.ceil(vlen(it["t"]) / per))
            fs = size if lv == 0 else size - 1.2
            total += lines * fs * 1.30 + (7 if lv == 0 else 3)
        return total

    def _bullets(self, s, items, x, y, w, h, size=None):
        if not items:
            return
        if size is None:
            size, ratio = self._fit(items, w, h)
            if ratio > 1.0:
                OVERFLOW.append((self.n, self._cur_title, round(ratio, 2)))
        tf = tbox(s, x, y, w, h)
        first = True
        for it in items:
            lv = it.get("lv", 0) or 0
            bold = bool(it.get("b"))
            txt = it["t"]
            mark = {0: "▪  ", 1: "–  ", 2: "·  "}[min(lv, 2)]
            fs = size if lv == 0 else size - 1.2
            col = INK if lv == 0 else MUTED if lv >= 2 else RGBColor(0x2C, 0x3A, 0x4B)
            if bold and lv == 0:
                col = ACCENT_D
            para(tf, first, mark + txt, fs, bold, col,
                 space_before=(6 if lv == 0 and not first else 2),
                 space_after=(1 if lv else 2),
                 indent=int(lv * 0.30 * 914400))
            first = False

    def _callout(self, s, text, y=None):
        if not text:
            return
        h = Inches(0.62)
        yy = y if y is not None else (BODY_BOT - h + Inches(0.02))
        box = rect(s, ML, yy, BODY_W, h, fill=AMBER_BG, line=RGBColor(0xF0, 0xD2, 0x9B))
        rect(s, ML, yy, Inches(0.07), h, fill=AMBER)
        tf = box.text_frame
        tf.margin_left = Inches(0.22)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "핵심   " + text
        style(r, 14 if vlen(text) < 110 else 12.5, True, AMBER)

    # ------------------------------------------------------------ 레이아웃들
    def bullets(self, title, items, subtitle=None, callout=None, note=None):
        s = self._new()
        top = self._chrome(s, title, subtitle)
        h = BODY_BOT - top - (Inches(0.78) if callout else 0)
        self._bullets(s, items, ML, top, BODY_W, h)
        self._callout(s, callout)
        self._note(s, note)
        return s

    def twocol(self, title, left, right, heads=None, subtitle=None, callout=None, note=None):
        s = self._new()
        top = self._chrome(s, title, subtitle)
        h = BODY_BOT - top - (Inches(0.78) if callout else 0)
        gap = Inches(0.34)
        cw = (BODY_W - gap) / 2
        hh = Inches(0.42) if heads else 0
        for i, (col, x) in enumerate(((left, ML), (right, ML + cw + gap))):
            if heads:
                b = rect(s, x, top, cw, Inches(0.36), fill=ACCENT if i == 0 else RGBColor(0x4B, 0x5A, 0x6B))
                label(b, heads[i], 13.5, True, WHITE)
            rect(s, x, top + hh, cw, h - hh, fill=RGBColor(0xFA, 0xFC, 0xFE), line=LINE)
            self._bullets(s, col, x + Inches(0.16), top + hh + Inches(0.14),
                          cw - Inches(0.32), h - hh - Inches(0.28))
        self._callout(s, callout)
        self._note(s, note)
        return s

    def code(self, title, lines, caption=None, items=None, subtitle=None, callout=None, note=None):
        s = self._new()
        top = self._chrome(s, title, subtitle)
        avail = BODY_BOT - top - (Inches(0.78) if callout else 0)
        bh = avail if not items else avail * 0.62
        if caption:
            tf = tbox(s, ML, top, BODY_W, Inches(0.26))
            para(tf, True, caption, 12.5, True, MUTED)
            top += Inches(0.32)
            bh -= Inches(0.32)
        cs = code_size(lines, bh)
        if code_rows(cs, bh) < len(lines):
            OVERFLOW.append((self.n, self._cur_title + " [코드]",
                             round(len(lines) / max(code_rows(cs, bh), 1), 2)))
        box = rect(s, ML, top, BODY_W, bh, fill=CODE_BG)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.22)
        tf.margin_top = Inches(0.14)
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.24
            p.space_after = 0
            r = p.add_run()
            r.text = ln if ln.strip() else " "
            col = RGBColor(0x7D, 0x91, 0xA8) if ln.strip().startswith(("//", "#")) else CODE_FG
            style(r, cs, False, col, mono=True)
        if items:
            self._bullets(s, items, ML, top + bh + Inches(0.18), BODY_W,
                          avail - bh - Inches(0.24))
        self._callout(s, callout)
        self._note(s, note)
        return s

    def table(self, title, headers, rows, subtitle=None, callout=None, note=None):
        s = self._new()
        top = self._chrome(s, title, subtitle)
        avail = BODY_BOT - top - (Inches(0.78) if callout else 0)
        nr, nc = len(rows) + 1, len(headers)
        gt = s.shapes.add_table(nr, nc, E(ML), E(top), E(BODY_W), E(min(avail, Inches(0.42) * nr)))
        tbl = gt.table
        tbl.first_row = True
        widths = self._colwidths(headers, rows)
        for j, frac in enumerate(widths):
            tbl.columns[j].width = Emu(E(BODY_W * frac))
        maxlen = max([vlen(h) for h in headers] +
                     [vlen(c) for r in rows for c in r] + [1])
        fs = 13 if maxlen <= 26 else (11.5 if maxlen <= 44 else (10 if maxlen <= 70 else 9))
        if nr > 12:
            fs = min(fs, 10)
        if nr > 16:
            fs = min(fs, 8.5)
        for j, htxt in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = htxt
            style(r, fs, True, WHITE)
        for i, row in enumerate(rows, start=1):
            for j in range(nc):
                cell = tbl.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF6, 0xF9, 0xFC)
                cell.margin_left = cell.margin_right = Inches(0.08)
                cell.margin_top = cell.margin_bottom = Inches(0.03)
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = row[j] if j < len(row) else ""
                style(r, fs, False, INK)
        self._callout(s, callout)
        self._note(s, note)
        return s

    @staticmethod
    def _colwidths(headers, rows):
        nc = len(headers)
        w = []
        for j in range(nc):
            m = vlen(headers[j])
            for r in rows:
                if j < len(r):
                    m = max(m, min(vlen(r[j]), 90))
            w.append(max(m, 6))
        tot = sum(w)
        frac = [x / tot for x in w]
        lo = 0.10
        frac = [max(lo, f) for f in frac]
        t = sum(frac)
        return [f / t for f in frac]

    def _note(self, s, note):
        if note:
            s.notes_slide.notes_text_frame.text = note

    # ------------------------------------------------------------ 논문 그림
    def diagram(self, title, img, items=None, subtitle=None, callout=None,
                note=None, source=None):
        """직접 그린 도해를 <객체>로 넣는다.

        그림은 이미 여백이 잘린 상태라, 그림 도형의 크기 = 선택 박스 크기가 된다.
        테두리를 그리지 않는다(도해가 자기 테두리를 갖는다).
        """
        return self.figure(title, img, source, items, subtitle, callout, note,
                           border=False)

    def figure(self, title, img, source=None, items=None, subtitle=None,
               callout=None, note=None, force=None, border=True):
        """논문에서 잘라낸 그림/표를 넣는다.

        가로로 긴 그림은 위에 꽉 채우고 설명을 아래에, 그렇지 않으면 왼쪽에 그림·오른쪽에 설명.
        """
        from PIL import Image as _Img
        s = self._new()
        top = self._chrome(s, title, subtitle)
        avail_h = BODY_BOT - top - (Inches(0.78) if callout else 0)
        iw, ih = _Img.open(img).size
        ar = iw / ih
        mode = force or ("top" if (ar >= 2.55 or not items) else "left")
        cap_h = Inches(0.24) if source else 0

        if mode == "top":
            box_w, box_h = BODY_W, avail_h * (0.62 if items else 1.0) - cap_h
            w = min(box_w, box_h * ar)
            h = w / ar
            x = ML + (BODY_W - w) / 2
            s.shapes.add_picture(img, E(x), E(top), E(w), E(h))
            if border:
                rect(s, x, top, w, h, fill=None, line=LINE, lw=0.75)
            if source:
                tf = tbox(s, ML, top + h + Inches(0.03), BODY_W, cap_h)
                para(tf, True, source, 10, False, MUTED, align=PP_ALIGN.CENTER)
            if items:
                iy = top + h + cap_h + Inches(0.14)
                self._bullets(s, items, ML, iy, BODY_W, BODY_BOT - iy -
                              (Inches(0.78) if callout else 0))
        else:
            gap = Inches(0.30)
            lw_ = BODY_W * 0.54
            box_h = avail_h - cap_h
            w = min(lw_, box_h * ar)
            h = w / ar
            s.shapes.add_picture(img, E(ML), E(top), E(w), E(h))
            if border:
                rect(s, ML, top, w, h, fill=None, line=LINE, lw=0.75)
            if source:
                tf = tbox(s, ML, top + h + Inches(0.03), w, cap_h)
                para(tf, True, source, 10, False, MUTED)
            rx = ML + lw_ + gap
            self._bullets(s, items or [], rx, top, SW - MR - rx, avail_h)
        self._callout(s, callout)
        self._note(s, note)
        return s

    def _bullets_cols(self, s, items, x, y, w, h, cols=2, gap=Inches(0.34), size=None):
        """설명을 여러 단으로 나눠 넣는다. 그림 아래 얕은 띠에 쓴다."""
        if not items:
            return
        if cols <= 1 or len(items) < 3:
            return self._bullets(s, items, x, y, w, h, size=size)
        # lv==0 을 기준으로 끊어야 하위 항목이 부모와 떨어지지 않는다
        groups, cur = [], []
        for it in items:
            if (it.get("lv", 0) or 0) == 0 and cur:
                groups.append(cur)
                cur = []
            cur.append(it)
        if cur:
            groups.append(cur)
        weight = [sum(max(1, vlen(i["t"]) / 46.0) for i in g) for g in groups]
        target = sum(weight) / cols
        cw = (w - gap * (cols - 1)) / cols
        col, acc = [], 0.0
        placed = 0
        size = size or self._fit(items, int(cw), int(h))[0]
        worst = 0.0
        chunks = []
        for g, wt in zip(groups, weight):
            if acc + wt > target * 1.06 and col and placed < cols - 1:
                chunks.append(col)
                placed += 1
                col, acc = [], 0.0
            col.extend(g)
            acc += wt
        if col:
            chunks.append(col)
        for i, ch in enumerate(chunks):
            self._bullets(s, ch, E(x + i * (cw + gap)), y, E(cw), h, size=size)
            worst = max(worst, self._fit(ch, int(cw), int(h), sizes=(size,))[1])
        if worst > 1.0:
            OVERFLOW.append((self.n, self._cur_title, round(worst, 2)))

    def diagram_svg(self, title, svg, items=None, subtitle=None, callout=None,
                    note=None, dia_h=Inches(4.60), cols=2, name=None):
        """SVG 도해를 **네이티브 도형 그룹**으로 넣는다(PNG 아님).

        상자·화살표·글자가 각각 PowerPoint 도형이라 그대로 편집할 수 있다.
        그림은 위에 폭 가득, 설명은 아래 띠에 여러 단으로 놓는다 —
        옆에 놓으면 그림 폭이 반으로 줄어 안쪽 글자가 5pt 까지 작아진다.
        """
        import sys as _sys
        import os as _os
        _sys.path.insert(0, _os.path.dirname(_os.path.abspath(__file__)))
        from svg2shapes import add_svg_shapes

        s = self._new()
        top = self._chrome(s, title, subtitle, tight=True)
        avail = BODY_BOT - top - (Inches(0.70) if callout else 0)
        # 설명이 실제로 차지할 높이만큼만 떼어 주고 나머지는 전부 그림에 준다.
        # 고정 비율로 나누면 설명이 짧은 장에서 그림이 공연히 작아진다.
        gap = Inches(0.34)
        cw = (BODY_W - gap * (cols - 1)) / cols
        if items:
            need = Inches(self._needed_pt(items, int(cw), float(self.DIA_ITEM_PT)) / cols / 72.0)
            need = max(Inches(0.80), need + Inches(0.10))   # 12pt 가 다 들어갈 만큼 뗀다 — 위 캡 없음
        else:
            need = 0
        dh = max(Inches(3.10), min(dia_h, avail - need - Inches(0.10)))
        # 그림은 예약 박스의 **위쪽**에 붙인다(valign='top'). 가로로 넓적한 도해는 폭 맞춤 뒤
        # 세로가 남는데, 세로 중앙에 두면 그림 위·아래에 빈 띠가 생긴다(L-25). 설명은
        # 예약 높이(dh)가 아니라 **실제 그림 높이(gh)** 바로 아래부터 시작해 남는 세로를 다 쓴다.
        g, gw, gh, n, minpt = add_svg_shapes(s.shapes, svg, E(ML), E(top),
                                             width=E(BODY_W), height=E(dh), valign="top")
        g.name = name or "도해"
        if items:
            by = top + min(dh, gh) + Inches(0.12)
            self._bullets_cols(s, items, ML, by, BODY_W, BODY_BOT - by -
                               (Inches(0.70) if callout else 0), cols=cols, gap=gap,
                               size=self.DIA_ITEM_PT)
        self._callout(s, callout)
        self._note(s, note)
        return s, n, minpt, (gw, gh)

    def save(self, path):
        self.prs.save(path)
        return path
