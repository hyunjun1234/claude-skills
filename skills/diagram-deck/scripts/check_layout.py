#!/usr/bin/env python3
"""슬라이드 레이아웃 점검 — 글자 넘침·도형 겹침·경계 이탈.

python-pptx 는 텍스트를 실제로 조판하지 않으므로, 글자 폭/높이를 추정해서 검사한다.
추정이라 완벽하진 않지만 '명백히 넘치는 것'은 잡는다.
"""
import re
import sys
from collections import defaultdict
from difflib import SequenceMatcher

from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0
PT_EMU = 12700.0


def vlen(s):
    return sum(2 if ord(c) > 0x2E80 else 1 for c in s)


def _real_width(seg, fs, name, bold, ea=None):
    """실제 글꼴로 잰 폭(pt). 못 재면 None.

    반각 어림짐작은 영문 라벨에서 10% 넘게 틀려 헛경고를 만든다.
    도해 안 라벨은 폭에 딱 맞춰 상자를 잡으므로 정확히 재야 한다.
    한글은 `a:ea` 글꼴로 그려지므로 그 이름으로 재야 한다.
    """
    try:
        import os
        import sys as _s
        _s.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        from svg2shapes import measure
        return measure(seg, fs, name or "맑은 고딕", bold, ea)
    except Exception:
        return None


def _ea_of(p):
    """문단 첫 run 의 East-Asian 타이프페이스."""
    from pptx.oxml.ns import qn
    for r in p.runs:
        rPr = r._r.find(qn("a:rPr"))
        if rPr is None:
            continue
        e = rPr.find(qn("a:ea"))
        if e is not None and e.get("typeface"):
            return e.get("typeface")
    return None


def est_size(shape):
    """(필요한 높이 pt, 가장 긴 줄의 폭 pt) 추정."""
    tf = shape.text_frame
    w_emu = shape.width or 0
    ml = tf.margin_left or 0
    mr = tf.margin_right or 0
    w_pt = max((w_emu - ml - mr) / PT_EMU, 10)
    total_h = 0.0
    max_w = 0.0
    for p in tf.paragraphs:
        sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
        fs = max(sizes) if sizes else 18.0
        mono = any((r.font.name or "") == "Consolas" for r in p.runs)
        per = fs * (0.55 if mono else 0.5)      # 반각 1칸의 폭
        fname = next((r.font.name for r in p.runs if r.font.name), None)
        fbold = any(bool(r.font.bold) for r in p.runs)
        fea = _ea_of(p)
        # p.text 는 <a:br/> 를 '\v' 로 준다 — 강제 줄바꿈으로 센다
        segs = (p.text or "").split("\v")
        lines = 0
        for seg in segs:
            line_w = _real_width(seg, fs, fname, fbold, fea) if seg else 0
            if line_w is None:
                line_w = vlen(seg) * per
            max_w = max(max_w, line_w)
            lines += max(1, int(line_w / w_pt) + (1 if line_w % w_pt else 0)) if seg else 1
        ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.2
        sb = p.space_before.pt if p.space_before is not None else 0
        sa = p.space_after.pt if p.space_after is not None else 0
        total_h += lines * fs * ls + sb + sa
    mt = (tf.margin_top or 0) / PT_EMU
    mb = (tf.margin_bottom or 0) / PT_EMU
    return total_h + mt + mb, max_w


def boxes(slide):
    """그룹 **안쪽까지** 훑는다.

    도해를 네이티브 도형 그룹으로 넣으면 슬라이드 최상위에는 그룹 하나만 보인다.
    그룹을 안 열면 도해 슬라이드는 사실상 무검사가 된다.
    """
    out = []

    def walk(container, depth=0):
        for sh in container:
            if sh.shape_type == 6 and depth < 4:        # GROUP
                walk(sh.shapes, depth + 1)
                continue
            if sh.left is None or sh.top is None:
                continue
            txt = sh.text_frame.text.strip() if sh.has_text_frame else ""
            out.append({
                "sh": sh, "txt": txt, "in_group": depth > 0,
                "l": sh.left, "t": sh.top,
                "r": sh.left + (sh.width or 0), "b": sh.top + (sh.height or 0),
            })
    walk(slide.shapes)
    return out


_PUNCT = re.compile(r"[\s'\"“”‘’·,.…—\-()\[\]:;!?{}=+·]")


def _norm(s):
    return _PUNCT.sub("", s)


def _maxpt(sh):
    m = 0.0
    if not sh.has_text_frame:
        return m
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                m = max(m, r.font.size.pt)
    return m


def text_heavy(bs, long_chars=12, long_ratio=0.75):
    """도해 그룹이 "그림"이 아니라 "글 상자 나열"인지 본다 (L-24).

    글 상자 덱의 지문: 그룹 안 글자 도형 중 12자 이상 문장이 3/4 이상이면서,
    선·격자·화살표 같은 비글자 도형이 글자 도형보다 적다.
    실측(2026-08-14): 글 상자 덱 5장 중 4장 검출, 그림 덱 4종 15장 중 1장(경계) 검출.
    """
    g = [b for b in bs if b["in_group"]]
    t = [b for b in g if b["txt"]]
    n = [b for b in g if not b["txt"]]
    if len(t) < 8:
        return None
    long = [b for b in t if len(b["txt"]) >= long_chars]
    ratio = len(long) / len(t)
    if ratio >= long_ratio and len(n) < len(t):
        return (len(t), len(n), round(ratio, 2))
    return None


_DASH = re.compile(r"\s[—–]\s")


def bullet_and_scale(prs, want_pt=12.0, width_ratio=0.97):
    """도해 슬라이드의 글머리 크기 고정(12pt)과 도해 축척 균일(폭 편차 ≤3%)을 본다 (L-30)."""
    bad_pt, widths = [], []
    for i, s in enumerate(prs.slides, 1):
        grp = [sh for sh in s.shapes if sh.shape_type == 6]
        if not grp:
            continue
        g = grp[0]
        widths.append((i, g.width))
        for sh in s.shapes:
            if sh.shape_type == 6 or not sh.has_text_frame:
                continue
            if sh.top < g.top + g.height:            # 도해 아래 = 글머리 영역
                continue
            if sh.top > prs.slide_height - 457200:   # 맨 아래 0.5in = 꼬리말 제외
                continue
            for pgh in sh.text_frame.paragraphs:
                for r in pgh.runs:
                    if r.font.size and abs(r.font.size.pt - want_pt) > 0.1:
                        bad_pt.append((i, round(r.font.size.pt, 1), r.text[:30]))
                        break
    scale_bad = []
    if len(widths) >= 2:
        mx = max(w for _, w in widths)
        for i, w in widths:
            if w < mx * width_ratio:
                scale_bad.append((i, round(w / mx, 3)))
    return bad_pt, scale_bad


_CIRC = "①②③④⑤⑥⑦⑧⑨"


def head_check(prs, want_pt=14.0):
    """도해 머리글 규약 (L-31): SVG 밖 슬라이드 층 · 14pt · 슬라이드마다 다른 번호."""
    in_svg, bad_pt, nums = [], [], []
    for i, s in enumerate(prs.slides, 1):
        for b in boxes(s):
            t = b["txt"].strip()
            if not t or t[0] not in _CIRC:
                continue
            first = t.split("\n")[0]
            if b["in_group"]:
                in_svg.append((i, first[:40]))          # 머리글이 SVG 안에 있음
                continue
            sh = b["sh"]
            r = next((r for pgh in sh.text_frame.paragraphs for r in pgh.runs if r.text.strip()), None)
            pt = r.font.size.pt if (r and r.font.size) else None
            if pt is None or abs(pt - want_pt) > 0.1:
                bad_pt.append((i, pt, first[:40]))
            nums.append((i, t[0]))
    dup = []
    seen = {}
    for i, c in nums:
        if c in seen:
            dup.append((i, c, seen[c]))
        seen.setdefault(c, i)
    return in_svg, bad_pt, dup


def accent_hues(prs, max_hues=2):
    """도해의 색 계열 수를 센다 (L-32). 무채색(잉크·회색) 제외 강조 색상군이
    max_hues(기본 2: 강조 1 + 하이라이트 1)를 넘으면 알록달록으로 보고한다."""
    import colorsys
    buckets = {}
    for i, s in enumerate(prs.slides, 1):
        for b in boxes(s):
            sh = b["sh"]
            cols = []
            try:
                if sh.fill.type is not None and str(sh.fill.type) == "MSO_FILL_TYPE.SOLID":
                    cols.append(sh.fill.fore_color.rgb)
            except Exception:
                pass
            try:
                cols.append(sh.line.color.rgb)
            except Exception:
                pass
            for c in cols:
                if c is None:
                    continue
                r, g, bl = c[0], c[1], c[2]
                if max(r, g, bl) - min(r, g, bl) < 28:      # 무채색
                    continue
                h = colorsys.rgb_to_hsv(r / 255, g / 255, bl / 255)[0]
                buckets.setdefault(round(h * 6) % 6, set()).add(f"#{c}")
    if len(buckets) > max_hues:
        return [(len(buckets), "; ".join(sorted(min(v) for v in buckets.values())))]
    return []


def label_fit(prs, tol=25400):
    """도해 안 라벨이 뒤에 깔린 상자를 넘는지 본다 (L-33).

    svg2shapes 는 라벨마다 독립 텍스트박스를 만들므로 기존 가로 넘침 검사가 못 본다.
    같은 그룹 안에서 라벨 중심을 포함하는 가장 작은 라벨급 상자(높이<0.9in, 폭<3.2in)를 찾아
    라벨 좌우가 상자를 tol(0.028in) 이상 벗어나면 보고한다."""
    hits = []
    IN = 914400
    for i, s in enumerate(prs.slides, 1):
        bs = boxes(s)
        texts = [b for b in bs if b["in_group"] and b["txt"] and "\n" not in b["txt"]]
        rects = [b for b in bs if b["in_group"] and not b["txt"]
                 and (b["b"] - b["t"]) < int(0.9 * IN) and (b["r"] - b["l"]) < int(3.2 * IN)
                 and (b["b"] - b["t"]) > int(0.10 * IN)]
        for t in texts:
            cx, cy = (t["l"] + t["r"]) // 2, (t["t"] + t["b"]) // 2
            cands = [r for r in rects if r["l"] <= cx <= r["r"] and r["t"] <= cy <= r["b"]]
            if not cands:
                continue
            r = min(cands, key=lambda q: (q["r"] - q["l"]) * (q["b"] - q["t"]))
            over = max(r["l"] - t["l"], t["r"] - r["r"])
            if over > tol:
                hits.append((i, "넘침", round(over / IN, 2), t["txt"][:30]))
            elif over > -int(0.08 * IN) and len(t["txt"]) >= 3 and (r["r"] - r["l"]) >= int(0.3 * IN):
                # 라벨이 상자에 거의 닿음 — 좌우 여유 0.08in 미만 (L-33; 글폭 추정이 낙관적이라 문턱 넉넉히).
                # 격자 칸의 한두 글자 숫자는 딱 맞는 게 정상이라 제외 (길이<3 or 상자 폭<0.3in).
                hits.append((i, "꽉참", round(-over / IN, 2), t["txt"][:30]))
    return hits


def dash_titles(prs):
    """'제목 — 덧붙임' 꼴을 잡는다 (L-26). 슬라이드 제목과 도해 첫 줄(①…) 이 대상."""
    hits = []
    for i, s in enumerate(prs.slides, 1):
        for b in boxes(s):
            t = b["txt"].strip()
            if not t or "\n" in t:
                continue
            is_title = (not b["in_group"]) and b["t"] < 914400   # 슬라이드 위 1in 안의 최상위 글자 = 제목
            is_head = b["in_group"] and t[:1] in "①②③④⑤⑥"
            if (is_title or is_head) and _DASH.search(t):
                hits.append((i, t[:60]))
    return hits


_SUMMARY = re.compile(r"(요약|정리|한 줄|summary|recap)", re.I)


def summary_last(prs):
    """마지막 슬라이드가 '요약/정리' 제목의 표·불릿이면 잡는다 (L-28)."""
    if len(prs.slides) < 2:
        return None
    s = prs.slides[-1]
    # deck.py 는 제목을 placeholder 가 아니라 맨 위 텍스트박스로 그린다 → 가장 위쪽 글자 도형을 제목으로 본다
    tb = [sh for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    if not tb:
        return None
    t = min(tb, key=lambda sh: sh.top).text_frame.text.strip().split("\n")[0]
    if _SUMMARY.search(t):
        return (len(prs.slides), t[:60])
    return None


def diagram_gap(prs, min_in=0.18, max_in=0.45):
    """도해 주변 간격을 잰다 (L-25). 너무 넓으면 빈 띠, 너무 좁으면 답답하다 —
    도해↔설명 간격이 [min_in, max_in] 인치 범위를 벗어나면 보고한다.
    머리글(원문자 시작 슬라이드 층 글자)↔도해 간격도 min_in 미만이면 보고."""
    hits = []
    IN = 914400
    for i, s in enumerate(prs.slides, 1):
        grp = [sh for sh in s.shapes if sh.shape_type == 6]
        if not grp:
            continue
        g = grp[0]
        below = [sh for sh in s.shapes if sh.shape_type != 6 and sh.has_text_frame and sh.text_frame.text.strip()
                 and sh.top >= g.top + g.height - 1000]
        if below:
            gap = (min(sh.top for sh in below) - (g.top + g.height)) / IN
            if gap > max_in:
                hits.append((i, "도해↔설명 넓음", round(gap, 2)))
            elif gap < min_in:
                hits.append((i, "도해↔설명 좁음", round(gap, 2)))
        heads = [sh for sh in s.shapes if sh.shape_type != 6 and sh.has_text_frame
                 and sh.text_frame.text.strip()[:1] in _CIRC and sh.top < g.top]
        for hd in heads:
            # 머리글 상자 높이가 아니라 글자 높이(14pt≈0.19in) 기준으로 잰다
            gap = (g.top - hd.top) / IN - 0.20
            if gap < min_in:
                hits.append((i, "머리글↔도해 좁음", round(gap, 2)))
    return hits


def dup_texts(bs, thr=0.80):
    """도해(그룹) 안 글자가 슬라이드 제목·핵심줄과 같은 말인지 본다.

    도해를 SVG 로 그릴 때 그림 안에 제목을 또 넣거나, 그림 아래 띠의 문장을
    콜아웃에 그대로 옮겨 적으면 한 슬라이드에 같은 말이 두 번 찍힌다(L-18).
    """
    inside = [b for b in bs if b["in_group"] and b["txt"]]
    if not inside:
        return []
    outside = [b for b in bs if not b["in_group"] and b["txt"]]
    refs = []
    titles = [(_maxpt(b["sh"]), b["txt"]) for b in outside]
    titles = [t for t in titles if t[0] >= 18]
    if titles:
        refs.append(("제목", max(titles, key=lambda t: t[0])[1]))
    for b in outside:
        if b["txt"].startswith("핵심"):
            refs.append(("핵심", b["txt"][2:].strip()))
    hits = []
    for label, ref in refs:
        rn = _norm(ref)
        if len(rn) < 8:
            continue
        for b in inside:
            tn = _norm(b["txt"])
            if len(tn) < 8:
                continue
            ratio = SequenceMatcher(None, rn, tn).ratio()
            if ratio >= thr:
                hits.append((label, ref[:36], b["txt"][:36], round(ratio, 2)))
                continue
            # 접두 포함: 도해 첫 줄이 "제목 + 덧붙임" 인 경우를 잡는다.
            # 유사도만 보면 덧붙임이 길수록 비율이 떨어져 빠져나간다(L-23).
            if len(rn) >= 8 and tn.startswith(rn):
                hits.append((label + "(접두)", ref[:36], b["txt"][:36], 1.0))
    return hits


ABSOLUTES = ("공통", "전부", "모두", "아무도", "뿐", "항상", "그대로",
             "언제나", "절대", "유일", "전혀", "무조건")


def absolutes(bs):
    """전칭 표현을 뽑는다 — 상세 슬라이드가 붙으면 가장 먼저 거짓이 되는 표현이다(L-19).

    자동 판정은 못 한다. 사람이 눈으로 확인하라고 뽑아 주는 것이다.
    """
    hits = []
    for b in bs:
        t = b["txt"]
        if not t:
            continue
        # 도해 안 라벨과 '핵심' 한 줄만 본다. 본문 불릿까지 넣으면 잡음이 너무 많아
        # 아무도 안 읽는다 — 짧게 단정하는 자리가 위험한 자리다.
        if not b["in_group"] and not t.startswith("핵심"):
            continue
        for w in ABSOLUTES:
            if w in t:
                hits.append((w, t[:60], "도해" if b["in_group"] else "핵심"))
                break
    return hits



def _is_solid(sh):
    """색이 칠해진 도형인가 (svg fill="none" 은 background 라 제외)."""
    try:
        return sh.fill.type == 1          # MSO_FILL.SOLID
    except Exception:
        return False


def seq_order(prs):
    """같은 접두어+번호 라벨(검증 1·검증 2…)이 읽기 순서로 놓였는지 (L-34).

    번호순으로 정렬해 중심 퍼짐이 큰 축을 진행 축으로 본다:
    세로 계열은 번호가 클수록 아래, 가로 계열은 오른쪽이어야 한다.
    (사다리 은유로 1번을 맨 아래 두면 읽기 순서가 뒤집힌다.)"""
    pat = re.compile(r"^(\D{1,12}?)\s*(\d{1,2})\s*[:.·]")
    hits = []
    for i, s in enumerate(prs.slides, 1):
        gr = {}
        for b in boxes(s):
            if not b["txt"] or "\n" in b["txt"]:
                continue
            m = pat.match(b["txt"])
            if not m or not m.group(1).strip():
                continue
            gr.setdefault(m.group(1).strip(), []).append((int(m.group(2)), b))
        for pre, it in gr.items():
            if len(it) < 2 or len({n for n, _ in it}) != len(it):
                continue
            it.sort(key=lambda q: q[0])
            xs = [(b["l"] + b["r"]) / 2 for _, b in it]
            ys = [(b["t"] + b["b"]) / 2 for _, b in it]
            if max(ys) - min(ys) >= max(xs) - min(xs):
                ok, ax = all(ys[k] < ys[k + 1] for k in range(len(ys) - 1)), "세로(위→아래)"
            else:
                ok, ax = all(xs[k] < xs[k + 1] for k in range(len(xs) - 1)), "가로(왼→오른)"
            if not ok:
                hits.append((i, pre, ax))
    return hits


def text_clash(prs, tol=int(0.03 * 914400)):
    """도해 글자가 **제 바탕이 아닌** 색칠 도형과 겹치는지 (L-33 확장).

    글자 중심이 그 도형 안이면 도형 위 라벨이므로 제외. 화살촉이 캡션을
    찌르거나 라벨이 옆 도형을 침범하면 양 축 모두 0.03in 넘게 겹친다."""
    hits = []
    for i, s in enumerate(prs.slides, 1):
        bs = boxes(s)
        texts = [b for b in bs if b["in_group"] and b["txt"] and "\n" not in b["txt"]]
        solids = [b for b in bs if b["in_group"] and not b["txt"] and _is_solid(b["sh"])]
        for t in texts:
            cx, cy = (t["l"] + t["r"]) // 2, (t["t"] + t["b"]) // 2
            for r in solids:
                if r["l"] <= cx <= r["r"] and r["t"] <= cy <= r["b"]:
                    continue
                ox = min(t["r"], r["r"]) - max(t["l"], r["l"])
                oy = min(t["b"], r["b"]) - max(t["t"], r["t"])
                if ox > tol and oy > tol:
                    hits.append((i, t["txt"][:26], round(ox / EMU_IN, 2), round(oy / EMU_IN, 2)))
                    break
    return hits


def covered_arrows(prs, frac=0.25):
    """앞서 그린 화살표·화살촉이 나중에 그린 색칠 도형에 덮이는지 (L-35).

    svg2shapes 는 문서 순서 = z순서다. 후보는 가는 도형(짧은 변 <0.06in,
    선)과 작은 도형(긴 변 ≤0.16in, 화살촉). 이후의 4배 이상 큰 솔리드가
    후보 넓이의 frac 이상을 덮으면 보고한다."""
    hits = []
    IN = EMU_IN
    for i, s in enumerate(prs.slides, 1):
        bs = boxes(s)
        for j, t in enumerate(bs):
            if t["txt"] or not t["in_group"]:
                continue
            w, h = t["r"] - t["l"], t["b"] - t["t"]
            if not (min(w, h) < int(0.06 * IN) or max(w, h) <= int(0.16 * IN)):
                continue
            l, tp, r, bt = t["l"], t["t"], t["r"], t["b"]
            if h < int(0.02 * IN):
                tp -= int(0.01 * IN); bt += int(0.01 * IN)
            if w < int(0.02 * IN):
                l -= int(0.01 * IN); r += int(0.01 * IN)
            area = max(1, (r - l) * (bt - tp))
            for c in bs[j + 1:]:
                if c["txt"] or not c["in_group"] or not _is_solid(c["sh"]):
                    continue
                if (c["r"] - c["l"]) * (c["b"] - c["t"]) < 4 * area:
                    continue
                ox = min(r, c["r"]) - max(l, c["l"])
                oy = min(bt, c["b"]) - max(tp, c["t"])
                if ox > 0 and oy > 0 and ox * oy > frac * area:
                    hits.append((i, ox * oy / area,
                                 round(t["l"] / IN, 2), round(t["t"] / IN, 2)))
                    break
    return hits


def box_slack(prs, min_slack=int(0.7 * 914400)):
    """왼쪽 정렬 라벨 상자의 오른쪽 빈 폭 과다 (L-36).

    상자 안 글자 묶음의 오른쪽 빈 폭이 0.7in 이상이면서 왼쪽 여백의 3배
    이상이면 — 상자 폭이 내용이 아니라 임의로 정해진 것이다. 폭을 내용에
    맞춰라. 컨테이너(높이 ≥1.2in 또는 폭 ≥6in)는 제외.
    **채움이 있는 도형만 본다** — 사선 리더선의 bounding box 를 상자로 오인하던
    오탐이 있었다(L-36 검사 수정)."""
    hits = []
    IN = EMU_IN
    for i, s in enumerate(prs.slides, 1):
        bs = boxes(s)
        texts = [b for b in bs if b["in_group"] and b["txt"] and "\n" not in b["txt"]]
        rects = [b for b in bs if b["in_group"] and not b["txt"] and _is_solid(b["sh"])
                 and int(0.10 * IN) < (b["b"] - b["t"]) < int(1.2 * IN)
                 and (b["r"] - b["l"]) < int(6.0 * IN)]
        for r in rects:
            ins = [t for t in texts
                   if r["l"] <= (t["l"] + t["r"]) // 2 <= r["r"]
                   and r["t"] <= (t["t"] + t["b"]) // 2 <= r["b"]]
            if not ins:
                continue
            sr = r["r"] - max(t["r"] for t in ins)
            sl = min(t["l"] for t in ins) - r["l"]
            if sr >= min_slack and sr >= 3 * max(sl, int(0.05 * IN)):
                hits.append((i, round(sr / IN, 2), ins[0]["txt"][:26]))
    return hits



def no_middot(prs):
    """PPT 안 모든 글에서 가운뎃점 금지 (L-37).

    나열은 쉼표, 짝·대안은 빗금(/), 곱은 × 로 쓴다. U+00B7 외에
    한글 가운뎃점(U+318D)·하이픈점(U+2027)도 같이 잡는다."""
    bad = ("\u00b7", "\u318d", "\u2027")
    hits = []
    for i, s in enumerate(prs.slides, 1):
        for b in boxes(s):
            if b["txt"] and any(ch in b["txt"] for ch in bad):
                hits.append((i, b["txt"][:40]))
    return hits


def overlap(a, b):
    ox = min(a["r"], b["r"]) - max(a["l"], b["l"])
    oy = min(a["b"], b["b"]) - max(a["t"], b["t"])
    if ox <= 0 or oy <= 0:
        return 0.0
    inter = ox * oy
    amin = min((a["r"] - a["l"]) * (a["b"] - a["t"]), (b["r"] - b["l"]) * (b["b"] - b["t"]))
    return inter / amin if amin else 0.0


def main(path):
    prs = Presentation(path)
    W, H = prs.slide_width, prs.slide_height
    over_h, over_w, collide, outside, dups, heavy = [], [], [], [], [], []
    for i, s in enumerate(prs.slides, 1):
        bs = boxes(s)
        for b in bs:
            sh = b["sh"]
            if b["l"] < -Emu(5000) or b["t"] < -Emu(5000) or b["r"] > W + Emu(20000) or b["b"] > H + Emu(20000):
                outside.append((i, b["txt"][:40], round(b["r"] / EMU_IN, 2), round(b["b"] / EMU_IN, 2)))
            if not b["txt"] or not sh.has_text_frame:
                continue
            need_h, need_w = est_size(sh)
            box_h = (sh.height or 0) / PT_EMU
            box_w = (sh.width or 0) / PT_EMU
            # 텍스트박스는 아래로 흘러넘쳐 아래 도형을 덮는다
            if need_h > box_h * 1.25 and need_h - box_h > 6:
                over_h.append((i, b["txt"][:44], round(need_h, 1), round(box_h, 1)))
            # 도형 라벨이 한 줄인데 폭을 넘으면 좌우로 삐져나온다
            if len(sh.text_frame.paragraphs) == 1 and need_w > box_w * 1.02 and box_h < 34:
                over_w.append((i, b["txt"][:44], round(need_w, 1), round(box_w, 1)))
        # 글자 있는 도형끼리 심하게 겹치는지
        tb = [b for b in bs if b["txt"]]
        for x in range(len(tb)):
            for y in range(x + 1, len(tb)):
                r = overlap(tb[x], tb[y])
                if r > 0.55:
                    collide.append((i, tb[x]["txt"][:26], tb[y]["txt"][:26], round(r, 2)))
        for h in dup_texts(bs):
            dups.append((i,) + h)
        th = text_heavy(bs)
        if th:
            heavy.append((i,) + th)

    dashes = dash_titles(prs)
    bad_pt, scale_bad = bullet_and_scale(prs)
    hd_svg, hd_pt, hd_dup = head_check(prs)
    hues = accent_hues(prs)
    lfit = label_fit(prs)
    summ = summary_last(prs)
    gaps = diagram_gap(prs)
    clash = text_clash(prs)
    cov = covered_arrows(prs)
    slack = box_slack(prs)
    seqs = seq_order(prs)
    dots = no_middot(prs)

    def rep(name, items, fmt):
        print(f"\n{'✅' if not items else '⚠️ '} {name}: {len(items)}건")
        for it in items[:18]:
            print("   " + fmt(it))
        if len(items) > 18:
            print(f"   … 외 {len(items)-18}건")

    print(f"슬라이드 {len(prs.slides._sldIdLst)}장 검사")
    rep("슬라이드 경계 이탈", outside, lambda x: f"p{x[0]:3d} R{x[2]} B{x[3]}  '{x[1]}'")
    rep("세로 넘침(아래 도형을 덮을 수 있음)", over_h,
        lambda x: f"p{x[0]:3d} 필요 {x[2]}pt / 박스 {x[3]}pt  '{x[1]}'")
    rep("가로 넘침(라벨이 도형 밖으로)", over_w,
        lambda x: f"p{x[0]:3d} 필요 {x[2]}pt / 박스 {x[3]}pt  '{x[1]}'")
    rep("글자 있는 도형끼리 겹침", collide,
        lambda x: f"p{x[0]:3d} {x[3]:.0%}  '{x[1]}' ↔ '{x[2]}'")
    rep("도해 안 글자가 슬라이드 글자와 중복", dups,
        lambda x: f"p{x[0]:3d} {x[4]:.0%} {x[1]}  '{x[2]}' ↔ 도해 '{x[3]}'")
    rep("도해 라벨이 상자를 넘거나 꽉 참 (L-33)", lfit, lambda x: f"p{x[0]:3d}  {x[1]} 여유 {x[2]}in  '{x[3]}'")
    rep("도해 글자가 이웃 도형을 침범 (L-33)", clash, lambda x: f"p{x[0]:3d}  겹침 {x[2]}×{x[3]}in  '{x[1]}'")
    rep("화살표·화살촉이 나중 도형에 덮임 (L-35)", cov, lambda x: f"p{x[0]:3d}  {x[1]:.0%} 덮임 @({x[2]},{x[3]})in")
    rep("상자 오른쪽 빈 폭 과다 (L-36)", slack, lambda x: f"p{x[0]:3d}  빈 폭 {x[1]}in  '{x[2]}'")
    rep("가운뎃점 '·' 사용 (L-37)", dots, lambda x: f"p{x[0]:3d}  '{x[1]}'")
    rep("순번 라벨이 읽기 순서 역행 (L-34)", seqs, lambda x: f"p{x[0]:3d}  '{x[1]} n' 계열 {x[2]} 역행")
    rep("색 계열 과다 — 무채 제외 3군 이상 (L-32)", hues, lambda x: f"강조 색상군 {x[0]}개: {x[1]}")
    rep("도해 머리글이 SVG 안에 있음 (L-31)", hd_svg, lambda x: f"p{x[0]:3d}  '{x[1]}'")
    rep("도해 머리글 크기가 14pt 아님 (L-31)", hd_pt, lambda x: f"p{x[0]:3d}  {x[1]}pt  '{x[2]}'")
    rep("도해 머리글 번호 중복 (L-31)", hd_dup, lambda x: f"p{x[0]:3d}  '{x[1]}' 이 p{x[2]} 와 중복")
    rep("도해 글머리 크기가 12pt 아님 (L-30)", bad_pt, lambda x: f"p{x[0]:3d}  {x[1]}pt  '{x[2]}'")
    rep("도해 축척 불일치(폭 편차 >3%) (L-30)", scale_bad, lambda x: f"p{x[0]:3d}  최대폭 대비 {x[1]:.0%}")
    rep("제목·도해 머리글의 '— 덧붙임' 꼴 (L-26)", dashes, lambda x: f"p{x[0]:3d}  '{x[1]}'")
    rep("마지막 슬라이드가 요약·정리 (L-28)", [summ] if summ else [], lambda x: f"p{x[0]:3d}  '{x[1]}'")
    rep("도해 주변 간격 범위 밖 (L-25, 0.18~0.45in)", gaps, lambda x: f"p{x[0]:3d}  {x[1]} {x[2]}in")
    rep("도해가 글 상자 나열(장문 라벨 ≥75% · 비글자 도형 부족)", heavy,
        lambda x: f"p{x[0]:3d} 글자도형 {x[1]} · 비글자 {x[2]} · 장문비율 {x[3]:.0%} — 구조를 그림으로 (L-24)")
    # text_heavy 는 경고다 — 실패 사유로 세지 않는다. 기존 덱을 일괄 실패시키지 않으면서
    # 새 덱을 만들 때 눈에 띄게 하는 것이 목적이다. 걸리면 그림으로 다시 그릴지 판단하라(L-24).
    bad = len(outside) + len(over_h) + len(over_w) + len(collide) + len(dups) + len(dashes) + (1 if summ else 0) + len(gaps) + len(bad_pt) + len(scale_bad) + len(hd_svg) + len(hd_pt) + len(hd_dup) + len(hues) + len(lfit) + len(clash) + len(cov) + len(slack) + len(seqs) + len(dots)
    print(f"\n{'✅ 레이아웃 문제 없음' if bad == 0 else f'총 {bad}건 확인 필요'}")
    return 1 if bad else 0


def report_absolutes(path):
    """검사가 아니라 검토 보조다. 전칭 표현이 있는 슬라이드를 모아 보여 준다."""
    prs = Presentation(path)
    n = 0
    for i, s_ in enumerate(prs.slides, 1):
        hits = absolutes(boxes(s_))
        if not hits:
            continue
        print(f"\np{i:3d}")
        for w, t, where in hits:
            n += 1
            print(f"   [{w}] ({where}) {t}")
    print(f"\n전칭 표현 {n}건 — 상세 슬라이드를 새로 넣었다면 이 문장들을 다시 읽어라(L-19).")
    return 0


if __name__ == "__main__":
    if "--absolutes" in sys.argv:
        sys.exit(report_absolutes(sys.argv[1]))
    sys.exit(main(sys.argv[1]))
