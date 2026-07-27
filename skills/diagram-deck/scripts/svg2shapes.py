#!/usr/bin/env python3
"""SVG 도해를 PowerPoint **네이티브 도형**으로 변환한다.

PNG 로 굽지 않는다. 상자·화살표·글자가 각각 편집 가능한 도형으로 들어가고,
전체는 그룹 하나로 묶여 옮기기 쉽다(그룹 안으로 들어가면 개별 편집).

지원 요소: svg, g(스타일 상속만), rect, circle, line, polygon, polyline, text
지원 속성: fill, stroke, stroke-width, stroke-dasharray, rx,
           font-size, font-weight, font-family, text-anchor
transform / tspan / style / class 는 지원하지 않는다 — SVG_RULES.md 가 금지한 것들이다.
"""
from __future__ import annotations

import re
import subprocess
from xml.etree import ElementTree as ET

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

EMU_PER_IN = 914400
INHERIT = ("fill", "stroke", "stroke-width", "stroke-dasharray",
           "font-size", "font-weight", "font-family", "text-anchor")

# SVG 글꼴 이름 → PPTX 에 적을 글꼴 이름.
# 폭은 그대로 로컬 Noto 로 재고(그래야 정확하다), 파일에는 어느 컴퓨터에나 있는 이름을 적는다.
# Noto Sans CJK KR 를 그대로 적으면 맥·윈도에서 대체 글꼴로 바뀌어 글자가 어긋난다.
FONT_MAP = {
    "Noto Sans CJK KR": "맑은 고딕",
    "Noto Sans Mono CJK KR": "Consolas",
}

# 출력 글꼴 → (여기서 폭을 대신 잴 글꼴, 실제 advance ÷ 그 글꼴의 advance).
# Consolas 는 리눅스에 없고 fc-match 가 DejaVu Sans Mono 로 대체해 버린다.
# 실제 Consolas 의 advance 는 0.55em, Noto Sans Mono CJK KR 은 0.50em → 1.10 배.
PROXY = {
    "Consolas": ("Noto Sans Mono CJK KR", 1.10),
    "맑은 고딕": ("Noto Sans CJK KR", 1.00),
}


def out_font(fam: str) -> str:
    return FONT_MAP.get(fam, fam)


def rendered_width(s: str, size: float, face: str, bold: bool) -> float:
    """그 글꼴로 실제 그렸을 때의 폭(pt) 추정. 검사기와 생성기가 같은 값을 봐야 한다."""
    fam, k = PROXY.get(face, (face, 1.0))
    return text_width(s, size, fam, bold) * k


def has_cjk(s: str) -> bool:
    return any(ord(c) > 0x2E80 for c in s)


def measure(s: str, size: float, latin: str, bold: bool, ea: str = None) -> float:
    """한글은 `a:ea` 글꼴로 그려진다.

    latin 이름(`Segoe UI` 등)으로 한글 폭을 재면 그 글꼴에 한글 글리프가 없어
    엉뚱하게 좁게 나온다 — 넘침 검사가 통째로 무력해진다.
    """
    return rendered_width(s, size, (ea or latin) if has_cjk(s) else latin, bold)


def face_of(s: str, latin: str, ea: str = None) -> str:
    return (ea or latin) if has_cjk(s) else latin


def size_for(face: str, size: float) -> float:
    """SVG 가 설계한 폭 그대로 나오도록 출력 글꼴의 크기를 보정한다.

    Consolas 로 그리면 같은 pt 에서 10% 넓어져 라벨이 상자를 삐져나온다.
    상자를 넓히는 대신 글자를 줄여야 그림의 배치가 안 깨진다.
    """
    return size / PROXY.get(face, (face, 1.0))[1]

# ── 글꼴 계측 ────────────────────────────────────────────────────────────────
# 글자 폭을 실제 글꼴로 잰다. 어림짐작(한글=1em, 영문=0.5em)은 오차가 커서
# 글자 상자가 그림 요소와 어긋난다.
_METRIC = {}
_REF = 1000  # 이 크기로 재고 비례로 줄인다


def _fontfile(family: str, bold: bool):
    key = f"{family}:{'bold' if bold else 'regular'}"
    try:
        out = subprocess.run(["fc-match", "-f", "%{file}|%{index}", key],
                             capture_output=True, text=True, timeout=5).stdout
        path, idx = out.split("|")
        return path, int(idx)
    except Exception:
        return None


def _metric(family: str, bold: bool):
    """(폭 측정 함수, ascent/em, descent/em) 을 준다."""
    key = (family, bold)
    if key in _METRIC:
        return _METRIC[key]
    fn = None
    hit = _fontfile(family, bold)
    if hit:
        try:
            from PIL import ImageFont
            fn = ImageFont.truetype(hit[0], _REF, index=hit[1])
        except Exception:
            fn = None
    if fn is None:                                   # 글꼴을 못 찾으면 어림짐작
        def width(s):
            return sum(_REF * (1.0 if ord(c) > 0x2000 else 0.52) for c in s)
        m = (width, 1.16, 0.32)
    else:
        asc, desc = fn.getmetrics()
        m = (lambda s: fn.getlength(s), asc / _REF, desc / _REF)
    _METRIC[key] = m
    return m


def text_width(s: str, size: float, family: str, bold: bool) -> float:
    w, _, _ = _metric(family, bold)
    return w(s) * size / _REF


# ── SVG 파싱 ────────────────────────────────────────────────────────────────
def _num(v, dflt=0.0):
    if v is None:
        return dflt
    try:
        return float(str(v).strip().rstrip("px"))
    except ValueError:
        return dflt


def _pts(s):
    v = [float(x) for x in re.split(r"[\s,]+", (s or "").strip()) if x]
    return list(zip(v[0::2], v[1::2]))


_NAMED = {"black": "000000", "white": "FFFFFF", "red": "FF0000",
          "blue": "0000FF", "green": "008000", "gray": "808080", "grey": "808080"}


def _color(v):
    if not v or v in ("none", "transparent"):
        return None
    v = v.strip()
    if v.startswith("#"):
        h = v[1:]
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        if len(h) == 6:
            return RGBColor.from_string(h.upper())
        return None
    h = _NAMED.get(v.lower())
    return RGBColor.from_string(h) if h else None


def parse(svg: str):
    """(prims, viewBox) 를 준다. prims 는 그리는 순서(= z 순서)대로."""
    root = ET.fromstring(svg)
    vb = [float(x) for x in re.split(r"[\s,]+", root.get("viewBox").strip())]
    prims = []

    def walk(el, inh):
        tag = el.tag.split("}")[-1]
        st = dict(inh)
        for k in INHERIT:
            if el.get(k) is not None:
                st[k] = el.get(k)
        if tag in ("svg", "g"):
            for ch in el:
                walk(ch, st)
        elif tag in ("rect", "circle", "line", "polygon", "polyline", "text"):
            prims.append((tag, el, st))
    walk(root, {})
    return prims, vb


def _style(el, st, key, dflt=None):
    v = el.get(key)
    return v if v is not None else st.get(key, dflt)


def _isbold(v):
    return str(v) in ("700", "800", "900", "bold", "bolder")


def _text_box(el, st):
    """글자 상자 (x, y, w, h). 경계 계산과 실제 삽입이 **같은 값**을 써야
    그룹 크기와 내용 경계가 정확히 일치한다."""
    fs = _num(_style(el, st, "font-size", 12), 12.0)
    fam = _style(el, st, "font-family", "Noto Sans CJK KR")
    bold = _isbold(_style(el, st, "font-weight", ""))
    anc = _style(el, st, "text-anchor", "start")
    s = (el.text or "").strip()
    tw = text_width(s, fs, fam, bold)
    _, asc, desc = _metric(fam, bold)
    # 편집 여유(pad)를 좌우 대칭으로 주면 정렬 기준점이 그만큼 밀린다.
    # start 는 왼쪽 끝, end 는 오른쪽 끝, middle 은 중심이 SVG 의 x 와 같아야 한다.
    pad = fs * 0.18
    w = tw + 2 * pad
    x, y = _num(el.get("x")), _num(el.get("y"))
    if anc == "middle":
        bx = x - w / 2
    elif anc == "end":
        bx = x - w
    else:
        bx = x
    return (bx, y - asc * fs, w, (asc + desc) * fs)


def _geom(tag, el, st):
    """(x0, y0, x1, y1) 도형이 차지하는 범위.

    선 두께는 넣지 않는다 — PowerPoint 의 도형 범위(ext)도 선을 경계선 위에
    걸쳐 그리고 범위에 포함하지 않는다. 여기서 포함하면 그룹 크기와 어긋난다.
    """
    if tag == "rect":
        x, y = _num(el.get("x")), _num(el.get("y"))
        return (x, y, x + _num(el.get("width")), y + _num(el.get("height")))
    if tag == "circle":
        cx, cy, r = _num(el.get("cx")), _num(el.get("cy")), _num(el.get("r"))
        return (cx - r, cy - r, cx + r, cy + r)
    if tag == "line":
        xs = [_num(el.get("x1")), _num(el.get("x2"))]
        ys = [_num(el.get("y1")), _num(el.get("y2"))]
        return (min(xs), min(ys), max(xs), max(ys))
    if tag in ("polygon", "polyline"):
        p = _pts(el.get("points"))
        xs = [q[0] for q in p] or [0.0]
        ys = [q[1] for q in p] or [0.0]
        return (min(xs), min(ys), max(xs), max(ys))
    x, y, w, h = _text_box(el, st)
    return (x, y, x + w, y + h)


def content_bbox(svg: str):
    """그림이 실제로 차지하는 영역. viewBox 가 아니라 이걸로 크기를 잡아야
    PPT 선택 박스가 그림에 딱 맞는다(LESSONS L-04)."""
    prims, vb = parse(svg)
    if not prims:
        return (0.0, 0.0, vb[2], vb[3])
    b = [_geom(t, e, s) for t, e, s in prims]
    return (min(q[0] for q in b), min(q[1] for q in b),
            max(q[2] for q in b), max(q[3] for q in b))


def font_sizes(svg: str):
    prims, _ = parse(svg)
    return [_num(_style(e, s, "font-size", 12), 12.0)
            for t, e, s in prims if t == "text" and (e.text or "").strip()]


def min_font(svg: str) -> float:
    fs = font_sizes(svg)
    return min(fs) if fs else 12.0


# ── PPTX 로 내보내기 ─────────────────────────────────────────────────────────
def _fill(shape, col):
    if col is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = col


def _stroke(shape, col, w_pt, dash):
    if col is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = col
    shape.line.width = Pt(max(0.25, w_pt))
    if dash:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def add_svg_shapes(shapes, svg: str, left, top, width=None, height=None):
    """SVG 를 네이티브 도형 그룹으로 슬라이드에 넣는다.

    left/top/width/height 는 EMU. width 나 height 중 하나만 줘도 비율로 맞춘다.
    반환: (group, w_emu, h_emu, 도형개수, 글자최소pt)
    """
    prims, _ = parse(svg)
    x0, y0, x1, y1 = content_bbox(svg)
    bw, bh = max(x1 - x0, 1e-6), max(y1 - y0, 1e-6)
    if width and height:
        k = min(width / bw, height / bh)
    elif width:
        k = width / bw
    else:
        k = height / bh
    ox = left + ((width - bw * k) / 2 if width else 0)
    oy = top + ((height - bh * k) / 2 if height else 0)

    def X(v):
        return Emu(int(round(ox + (v - x0) * k)))

    def Y(v):
        return Emu(int(round(oy + (v - y0) * k)))

    def L(v):                       # 길이(원점 보정 없음)
        return Emu(max(1, int(round(v * k))))

    pt = k / EMU_PER_IN * 72.0       # user unit -> pt
    g = shapes.add_group_shape()
    gs = g.shapes
    n = 0
    smallest = None

    for tag, el, st in prims:
        fill = _color(_style(el, st, "fill", "#000000"))
        stroke = _color(_style(el, st, "stroke"))
        sw = _num(_style(el, st, "stroke-width", 1), 1.0) * pt
        dash = bool(_style(el, st, "stroke-dasharray"))

        if tag == "rect":
            w, h = _num(el.get("width")), _num(el.get("height"))
            rx = _num(el.get("rx"))
            kind = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
            sp = gs.add_shape(kind, X(_num(el.get("x"))), Y(_num(el.get("y"))),
                              L(w), L(h))
            if rx > 0 and min(w, h) > 0:
                try:
                    sp.adjustments[0] = min(0.5, rx / min(w, h))
                except Exception:
                    pass
            sp.shadow.inherit = False
            _fill(sp, fill)
            _stroke(sp, stroke, sw, dash)
        elif tag == "circle":
            cx, cy, r = _num(el.get("cx")), _num(el.get("cy")), _num(el.get("r"))
            sp = gs.add_shape(MSO_SHAPE.OVAL, X(cx - r), Y(cy - r),
                              L(2 * r), L(2 * r))
            sp.shadow.inherit = False
            _fill(sp, fill)
            _stroke(sp, stroke, sw, dash)
        elif tag == "line":
            sp = gs.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  X(_num(el.get("x1"))), Y(_num(el.get("y1"))),
                                  X(_num(el.get("x2"))), Y(_num(el.get("y2"))))
            _stroke(sp, stroke or RGBColor(0, 0, 0), sw, dash)
        elif tag in ("polygon", "polyline"):
            p = _pts(el.get("points"))
            if len(p) < 2:
                continue
            fb = gs.build_freeform(X(p[0][0]), Y(p[0][1]), scale=1.0)
            fb.add_line_segments([(X(a), Y(b)) for a, b in p[1:]],
                                 close=(tag == "polygon"))
            sp = fb.convert_to_shape()
            sp.shadow.inherit = False
            _fill(sp, fill if tag == "polygon" else None)
            _stroke(sp, stroke, sw, dash)
        else:  # text
            s = (el.text or "").strip()
            if not s:
                continue
            fs = _num(_style(el, st, "font-size", 12), 12.0)
            fam = _style(el, st, "font-family", "Noto Sans CJK KR")
            bold = _isbold(_style(el, st, "font-weight", ""))
            anc = _style(el, st, "text-anchor", "start")
            bx, by, bw_, bh_ = _text_box(el, st)
            tb = gs.add_textbox(X(bx), Y(by), L(bw_), L(bh_))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p0 = tf.paragraphs[0]
            p0.alignment = {"middle": PP_ALIGN.CENTER,
                            "end": PP_ALIGN.RIGHT}.get(anc, PP_ALIGN.LEFT)
            r = p0.add_run()
            r.text = s
            face = out_font(fam)
            size_pt = round(size_for(face, fs * pt), 1)
            smallest = size_pt if smallest is None else min(smallest, size_pt)
            f = r.font
            f.size = Pt(size_pt)
            f.bold = bold
            f.name = face
            f.color.rgb = fill if fill is not None else RGBColor(0x11, 0x11, 0x11)
            _ea_cs(r, face)                   # latin -> ea -> cs (LESSONS L-01)
        n += 1

    w_emu, h_emu = fix_group_extents(g)
    return g, w_emu, h_emu, n, (smallest or 0.0)


def _ea_cs(run, face):
    """DrawingML CT_TextCharacterProperties 는 latin -> ea -> cs 순서를 강제한다."""
    from lxml import etree
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = etree.SubElement(rPr, qn("a:cs"))
    if latin is not None:
        latin.addnext(cs)
        latin.addnext(ea)
    ea.set("typeface", face)
    cs.set("typeface", face)


def fix_group_extents(group):
    """그룹 크기를 자식들의 실제 경계로 다시 잡는다.

    python-pptx 는 add_shape/add_connector 때만 자동 갱신하고
    build_freeform 으로 만든 도형은 반영하지 않는다. 그대로 두면
    그룹 선택 박스가 그림과 어긋난다(LESSONS L-04).
    """
    kids = list(group.shapes)
    if not kids:
        return 0, 0
    x0 = min(s.left for s in kids)
    y0 = min(s.top for s in kids)
    x1 = max(s.left + s.width for s in kids)
    y1 = max(s.top + s.height for s in kids)
    xfrm = group._element.find(qn("p:grpSpPr")).find(qn("a:xfrm"))
    for tag in ("a:off", "a:chOff"):
        e = xfrm.find(qn(tag))
        e.set("x", str(int(x0)))
        e.set("y", str(int(y0)))
    for tag in ("a:ext", "a:chExt"):
        e = xfrm.find(qn(tag))
        e.set("cx", str(int(x1 - x0)))
        e.set("cy", str(int(y1 - y0)))
    return int(x1 - x0), int(y1 - y0)


DEMO = ('<svg viewBox="0 0 400 160" xmlns="http://www.w3.org/2000/svg">'
        '<g font-family="Noto Sans CJK KR" font-size="13">'
        '<rect x="20" y="30" width="150" height="60" rx="8" fill="#e8f0fb" '
        'stroke="#0b5fc0" stroke-width="1.5"/>'
        '<text x="95" y="65" text-anchor="middle" font-weight="700">입력 타일</text>'
        '<line x1="175" y1="60" x2="235" y2="60" stroke="#5b6b7b" stroke-width="1.6"/>'
        '<polygon points="235,60 227,56 227,64" fill="#5b6b7b"/>'
        '<circle cx="300" cy="60" r="26" fill="#fdece9" stroke="#b4222a" stroke-width="1.4"/>'
        '<text x="300" y="64" text-anchor="middle">PE</text></g></svg>')

if __name__ == "__main__":
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    g, w, h, n, mn = add_svg_shapes(sl.shapes, DEMO, Emu(914400), Emu(914400),
                                    width=Emu(6 * EMU_PER_IN))
    print(f"도형 {n}개, 그룹 {w/EMU_PER_IN:.2f} x {h/EMU_PER_IN:.2f} in, 최소 글자 {mn}pt")
    prs.save("/tmp/svg2shapes_demo.pptx")
    print("→ /tmp/svg2shapes_demo.pptx")
