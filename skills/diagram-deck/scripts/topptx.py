#!/usr/bin/env python3
"""PPTX 변환 코어 — HTML 도, PDF 도 여기로 들어와 곧장 PPTX 가 된다.

HTML 은 PDF 파일을 거치지 않는다. WeasyPrint 가 낸 바이트를 메모리에서 바로 받아 쓴다.
(중간 PDF 를 굳이 남기고 싶으면 keep_pdf 로 따로 저장할 수 있다.)
"""
import os

import fitz
from pptx import Presentation
from pptx.oxml.ns import qn

import deck


def _pptx_from_pdf_bytes(pdf_bytes, out, dpi=200, notes=True, tmpdir=None):
    doc = fitz.open("pdf", pdf_bytes)
    pw, ph = doc[0].rect.width, doc[0].rect.height

    prs = Presentation()
    prs.slide_width = deck.E(pw * 12700)          # 1 pt = 12700 EMU
    prs.slide_height = deck.E(ph * 12700)
    blank = prs.slide_layouts[6]

    tmpdir = tmpdir or os.path.join(os.path.dirname(os.path.abspath(out)), "_topptx_png")
    os.makedirs(tmpdir, exist_ok=True)
    zoom = dpi / 72.0

    for i, page in enumerate(doc, 1):
        png = os.path.join(tmpdir, f"p{i:03d}.png")
        page.get_pixmap(matrix=fitz.Matrix(zoom, zoom)).save(png)
        s = prs.slides.add_slide(blank)
        # 슬라이드에는 그림 도형 하나만 둔다 — 손상 위험과 선택 박스 문제를 동시에 피한다
        s.shapes.add_picture(png, 0, 0, prs.slide_width, prs.slide_height)

        txt = page.get_text().strip()
        if notes and txt:
            tf = s.notes_slide.notes_text_frame
            tf.word_wrap = True
            deck._runs(tf.paragraphs[0], txt, 11, False, deck.INK, False)
        # 개요에 뜨도록 슬라이드 이름을 붙인다 (숨은 텍스트 상자를 쓰지 않는다)
        first = next((ln.strip() for ln in txt.split("\n") if len(ln.strip()) > 3), f"{i}쪽")
        csld = s._element.find(qn("p:cSld"))
        if csld is not None:
            csld.set("name", first[:60])

    prs.save(out)
    return {"pages": doc.page_count,
            "slide_in": (prs.slide_width / 914400, prs.slide_height / 914400),
            "page_pt": (pw, ph), "dpi": dpi,
            "bytes": os.path.getsize(out), "png_dir": tmpdir}


def html_to_pptx(html_path, out, dpi=200, notes=True, keep_pdf=None, base_url=None):
    """HTML → PPTX. 중간 PDF 파일을 만들지 않는다."""
    from weasyprint import HTML
    base = base_url or os.path.dirname(os.path.abspath(html_path))
    pdf_bytes = HTML(filename=html_path, base_url=base).write_pdf()
    if keep_pdf:
        with open(keep_pdf, "wb") as f:
            f.write(pdf_bytes)
    return _pptx_from_pdf_bytes(pdf_bytes, out, dpi, notes)


def pdf_to_pptx(pdf_path, out, dpi=200, notes=True):
    """PDF → PPTX."""
    with open(pdf_path, "rb") as f:
        return _pptx_from_pdf_bytes(f.read(), out, dpi, notes)


def report(info, src, out):
    print(f"{src} → {out}")
    print(f"  {info['pages']}쪽 → 슬라이드 {info['pages']}장, "
          f"{info['slide_in'][0]:.2f} x {info['slide_in'][1]:.2f} in "
          f"(원본 {info['page_pt'][0]:.0f}x{info['page_pt'][1]:.0f}pt), "
          f"{info['dpi']} DPI, {info['bytes']/1048576:.2f} MB")
