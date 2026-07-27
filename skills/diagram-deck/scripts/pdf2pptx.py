#!/usr/bin/env python3
"""PDF → PPTX 변환.

각 쪽을 고해상도 PNG 로 렌더해 슬라이드 하나에 정확히 채운다.
슬라이드 크기를 원본 쪽 크기와 같게 잡으므로 여백도 왜곡도 생기지 않는다.
쪽의 텍스트는 발표자 노트에 넣어 검색·복사가 되게 한다.

PowerPoint 손상 판정을 피하려고 지키는 것:
  - 좌표는 전부 정수 EMU (ST_Coordinate = xsd:long)
  - 글꼴은 latin -> ea -> cs 순서로 지정 (deck.style 재사용)
  - 줄바꿈은 <a:br/> (deck._runs 재사용)
"""
import os
import sys

import fitz
from pptx import Presentation
from pptx.oxml.ns import qn

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import deck  # noqa: E402

SRC = sys.argv[1]
OUT = sys.argv[2]
DPI = int(sys.argv[3]) if len(sys.argv) > 3 else 200

doc = fitz.open(SRC)
pw, ph = doc[0].rect.width, doc[0].rect.height          # pt
prs = Presentation()
prs.slide_width = deck.E(pw * 12700)                    # 1 pt = 12700 EMU
prs.slide_height = deck.E(ph * 12700)
blank = prs.slide_layouts[6]

tmpdir = os.path.join(os.path.dirname(os.path.abspath(OUT)), "_pdf2pptx_png")
os.makedirs(tmpdir, exist_ok=True)
zoom = DPI / 72.0

for i, page in enumerate(doc, 1):
    png = os.path.join(tmpdir, f"p{i:03d}.png")
    page.get_pixmap(matrix=fitz.Matrix(zoom, zoom)).save(png)
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(png, 0, 0, prs.slide_width, prs.slide_height)

    # 발표자 노트: 그 쪽의 텍스트 (검색·복사용)
    txt = page.get_text().strip()
    if txt:
        tf = s.notes_slide.notes_text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        deck._runs(p, txt, 11, False, deck.INK, False)

    # 슬라이드 이름을 붙인다 (p:cSld/@name). 숨은 텍스트 상자를 두지 않으므로
    # 슬라이드에는 그림 도형 하나만 남고, 그만큼 손상 위험이 줄어든다.
    first = next((ln.strip() for ln in txt.split("\n") if len(ln.strip()) > 3), f"{i}쪽")
    csld = s._element.find(qn("p:cSld"))
    if csld is not None:
        csld.set("name", first[:60])

prs.save(OUT)
print(f"{doc.page_count}쪽 → {OUT}")
print(f"슬라이드 {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} in "
      f"(원본 {pw:.0f}x{ph:.0f}pt), {DPI} DPI, {os.path.getsize(OUT)/1048576:.2f} MB")
