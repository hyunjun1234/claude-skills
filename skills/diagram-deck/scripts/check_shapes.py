#!/usr/bin/env python3
"""네이티브 도형 변환 검증 — PPTX 안의 도형을 SVG 로 되돌려 원본과 비교한다.

`svg2shapes.add_svg_shapes()` 로 넣은 그림이 정말 원본대로 들어갔는지 본다.
두 가지를 한다.

1. **되돌리기 렌더 비교** — pptx 의 도형 XML(위치·색·글자)을 다시 SVG 로 그려
   원본 SVG 와 같은 크기로 렌더한 뒤 픽셀 차이를 잰다. 요소가 빠지거나
   좌표·색이 틀리면 차이로 드러난다.
2. **그룹 경계 검사** — 그룹 크기가 자식들의 실제 경계와 같은지 본다
   (선택 박스가 그림보다 크면 안 된다, LESSONS L-04).

한계: PowerPoint 자체 렌더러가 아니라 이 검사기의 해석으로 그린 것이다.
글꼴 계측이 PowerPoint 와 완전히 같지는 않다. 좌표·색·누락은 잡지만
'PowerPoint 에서 글자가 상자를 넘는가' 는 못 잡는다.
"""
from __future__ import annotations

import io
import os
import sys

from pptx.oxml.ns import qn
from pptx.util import Emu

EMU_PER_IN = 914400
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _hex(c):
    return f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}" if c else "none"


def _solid_color(el):
    """spPr 또는 ln 요소에서 solidFill 색을 꺼낸다."""
    if el is None:
        return None
    if el.find(qn("a:noFill")) is not None:
        return None
    sf = el.find(qn("a:solidFill"))
    if sf is None:
        return None
    srgb = sf.find(qn("a:srgbClr"))
    if srgb is None:
        return None
    v = srgb.get("val")
    return (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def _xfrm(sp):
    x = sp.find(qn("a:xfrm")) if sp is not None else None
    if x is None:
        return None
    off, ext = x.find(qn("a:off")), x.find(qn("a:ext"))
    return (int(off.get("x")), int(off.get("y")),
            int(ext.get("cx")), int(ext.get("cy")),
            x.get("flipH") == "1", x.get("flipV") == "1")


def _flat(shapes, depth=0):
    """그룹을 펼쳐 실제로 그려지는 도형만 순서대로 준다."""
    out = []
    for sh in shapes:
        if sh.shape_type == 6 and depth < 4:
            out.extend(_flat(sh.shapes, depth + 1))
        else:
            out.append(sh)
    return out


def slide_to_svg(slide, slide_w, slide_h, dpi=110):
    """슬라이드 한 장을 통째로 SVG 로 되돌린다.

    LibreOffice 없이 배치를 눈으로 확인하려고 쓴다. 정확한 조판이 아니라
    '무엇이 어디에 있는가'를 보는 용도다.
    """
    return _emit(_flat(slide.shapes), 0, 0, slide_w, slide_h, dpi)


def group_to_svg(group, dpi=110):
    """그룹 안 도형들을 SVG 문자열로 되돌린다. (svg, (폭EMU, 높이EMU))"""
    kids = _flat(group.shapes)
    x0 = min(s.left for s in kids)
    y0 = min(s.top for s in kids)
    x1 = max(s.left + s.width for s in kids)
    y1 = max(s.top + s.height for s in kids)
    return _emit(kids, x0, y0, x1, y1, dpi), (x1 - x0, y1 - y0)


def _emit(kids, x0, y0, x1, y1, dpi=110):
    sc = dpi / EMU_PER_IN

    def P(v):
        return round((v) * sc, 2)

    W, H = P(x1 - x0), P(y1 - y0)
    out = [f'<svg viewBox="0 0 {W} {H}" width="{W}pt" height="{H}pt" '
           f'xmlns="http://www.w3.org/2000/svg">']

    for sh in kids:
        el = sh._element
        tag = el.tag.split("}")[-1]
        spPr = el.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln")) if spPr is not None else None
        fill = _solid_color(spPr)
        stroke = _solid_color(ln)
        lw = 1.0
        if ln is not None and ln.get("w"):
            lw = int(ln.get("w")) / 12700.0 * dpi / 72.0
        dash = ln is not None and ln.find(qn("a:prstDash")) is not None
        da = ' stroke-dasharray="4 3"' if dash else ""
        xf = _xfrm(spPr)
        if xf is None:                          # 자리표시자 등 — 그릴 위치가 없다
            continue
        L, T, Wd, Ht, fh, fv = xf
        L, T = L - x0, T - y0
        if tag == "pic":                        # 그림은 회색 상자로만 표시(배치 확인용)
            out.append(f'<rect x="{P(L)}" y="{P(T)}" width="{P(Wd)}" height="{P(Ht)}" '
                       f'fill="#e6ebf1" stroke="#9aa8b6" stroke-width="1"/>'
                       f'<text x="{P(L+Wd/2)}" y="{P(T+Ht/2)}" text-anchor="middle" '
                       f'font-family="Noto Sans CJK KR" font-size="11" fill="#5b6b7b">'
                       f'[그림]</text>')
            continue

        if tag == "cxnSp":                      # 직선
            ax, ay = (L + Wd, T) if fh else (L, T)
            bx, by = (L, T + Ht) if fh else (L + Wd, T + Ht)
            if fv:
                ay, by = by, ay
            out.append(f'<line x1="{P(ax)}" y1="{P(ay)}" x2="{P(bx)}" y2="{P(by)}" '
                       f'stroke="{_hex(stroke)}" stroke-width="{lw:.2f}"{da}/>')
            continue

        prst = spPr.find(qn("a:prstGeom"))
        cust = spPr.find(qn("a:custGeom"))
        if cust is not None:                    # 자유형(화살촉 등)
            pts = []
            for pt in cust.iter(qn("a:pt")):
                pts.append((L + int(pt.get("x")), T + int(pt.get("y"))))
            closed = cust.find(f".//{A}close") is not None
            d = " ".join(f"{P(a)},{P(b)}" for a, b in pts)
            t = "polygon" if closed else "polyline"
            out.append(f'<{t} points="{d}" fill="{_hex(fill)}" '
                       f'stroke="{_hex(stroke)}" stroke-width="{lw:.2f}"{da}/>')
            continue
        if prst is not None:
            kind = prst.get("prst")
            if kind == "ellipse":
                out.append(f'<ellipse cx="{P(L+Wd/2)}" cy="{P(T+Ht/2)}" '
                           f'rx="{P(Wd/2)}" ry="{P(Ht/2)}" fill="{_hex(fill)}" '
                           f'stroke="{_hex(stroke)}" stroke-width="{lw:.2f}"{da}/>')
                continue
            rx = 0.0
            if kind == "roundRect":
                gd = prst.find(qn("a:avLst"))
                v = 16667
                if gd is not None:
                    g0 = gd.find(qn("a:gd"))
                    if g0 is not None:
                        v = float(g0.get("fmla").split()[-1])
                rx = v / 100000.0 * min(Wd, Ht)
            txt = "".join(t.text or "" for t in el.iter(qn("a:t"))).strip()
            if txt or fill is not None or stroke is not None:
                out.append(f'<rect x="{P(L)}" y="{P(T)}" width="{P(Wd)}" '
                           f'height="{P(Ht)}" rx="{P(rx)}" fill="{_hex(fill)}" '
                           f'stroke="{_hex(stroke)}" stroke-width="{lw:.2f}"{da}/>')
            if not txt:
                continue

        # 글자 — 문단마다, 줄바꿈·자동 줄나눔까지 반영한다.
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        out.extend(_text_svg(sh, L, T, Wd, Ht, P, dpi))
    out.append("</svg>")
    return "\n".join(out)


def _text_svg(sh, L, T, Wd, Ht, P, dpi):
    """도형 하나의 글자를 SVG <text> 여러 개로 그린다."""
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from svg2shapes import PROXY, _metric, face_of, rendered_width

    el = sh._element
    bpr = next(el.iter(qn("a:bodyPr")), None)
    mL = int(bpr.get("lIns", 91440)) if bpr is not None else 91440
    mR = int(bpr.get("rIns", 91440)) if bpr is not None else 91440
    mT = int(bpr.get("tIns", 45720)) if bpr is not None else 45720
    wrap = (bpr.get("wrap") != "none") if bpr is not None else True
    anch = bpr.get("anchor") if bpr is not None else None
    iw = max(Wd - mL - mR, 1)
    iw_pt = iw / 12700.0

    lines = []                                   # (텍스트, fs, bold, fam, 색, 정렬, 들여쓰기, 위여백)
    for p in sh.text_frame.paragraphs:
        runs = p.runs
        if not runs:
            lines.append((None, 12.0, False, "맑은 고딕", None, "start", 0, 0))
            continue
        fs = max((r.font.size.pt for r in runs if r.font.size), default=18.0)
        bold = any(bool(r.font.bold) for r in runs)
        fam = next((r.font.name for r in runs if r.font.name), "맑은 고딕")
        ea = None
        for r in runs:                       # 한글은 a:ea 글꼴로 그려진다
            rPr = r._r.find(qn("a:rPr"))
            e = rPr.find(qn("a:ea")) if rPr is not None else None
            if e is not None and e.get("typeface"):
                ea = e.get("typeface")
                break
        col = next((r.font.color.rgb for r in runs
                    if r.font.color and r.font.color.type is not None), None)
        algn = {2: "middle", 3: "end"}.get(
            int(p.alignment) if p.alignment is not None else 1, "start")
        ind = (p._pPr.get("marL") if p._pPr is not None else None)
        ind = int(ind) if ind else 0
        sb = p.space_before.pt if p.space_before is not None else 0
        avail = iw_pt - ind / 12700.0
        for seg in (p.text or "").split("\v"):
            pfam = PROXY.get(face_of(seg, fam, ea), (face_of(seg, fam, ea), 1.0))[0]
            if not wrap:
                lines.append((seg, fs, bold, pfam, col, algn, ind, sb))
                sb = 0
                continue
            cur = ""
            for ch in seg:                        # 글자 단위 줄나눔(대충이지만 충분하다)
                if rendered_width(cur + ch, fs, face_of(cur + ch, fam, ea),
                                  bold) > avail and cur:
                    lines.append((cur, fs, bold, pfam, col, algn, ind, sb))
                    sb = 0
                    cur = ch
                else:
                    cur += ch
            lines.append((cur, fs, bold, pfam, col, algn, ind, sb))
            sb = 0

    total = sum((f * 1.30 + s) for _t, f, _b, _fa, _c, _a, _i, s in lines)
    y0 = T / EMU_PER_IN * dpi + mT / EMU_PER_IN * dpi
    if anch in ("ctr", "b"):
        box_h_px = Ht / EMU_PER_IN * dpi
        tot_px = total / 72.0 * dpi
        y0 = T / EMU_PER_IN * dpi + (box_h_px - tot_px) * (0.5 if anch == "ctr" else 1.0)
    out, y = [], y0
    for txt, fs, bold, fam, col, algn, ind, sb in lines:
        fs_px = fs / 72.0 * dpi
        y += sb / 72.0 * dpi
        _, asc, _d = _metric(fam, bold)
        if txt:
            x = L + mL + ind + (iw / 2 if algn == "middle" else (iw if algn == "end" else 0))
            esc = txt.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            out.append(f'<text x="{P(x)}" y="{round(y + asc*fs_px, 2)}" '
                       f'text-anchor="{algn}" font-family="{fam}" '
                       f'font-size="{fs_px:.2f}" font-weight="{700 if bold else 400}" '
                       f'fill="{_hex(col) if col else "#111111"}">{esc}</text>')
        y += fs_px * 1.30
    return out


def render_png(svg, px_w=900):
    """SVG 를 폭 px_w 픽셀로 렌더한다. 높이는 viewBox 비율대로."""
    import re

    import fitz
    from PIL import Image
    from weasyprint import HTML
    m = re.search(r'viewBox\s*=\s*["\']([\d.\-\s]+)["\']', svg)
    _, _, vw, vh = [float(v) for v in re.split(r"[\s,]+", m.group(1).strip())]
    html = ('<meta charset="utf-8"><style>'
            f'@page{{size:{vw}pt {vh}pt;margin:0}}'
            'html,body{margin:0;padding:0;background:#fff}svg{display:block}</style>'
            + re.sub(r"<svg\b", f'<svg width="{vw}pt" height="{vh}pt"', svg, count=1))
    pdf = HTML(string=html).write_pdf()
    z = px_w / vw
    pix = fitz.open("pdf", pdf)[0].get_pixmap(matrix=fitz.Matrix(z, z), alpha=False)
    return Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")


def crop_to_content(svg):
    """원본 SVG 의 viewBox 를 내용 경계로 바꾼다.

    되돌리기 렌더는 그룹(=내용 경계)만 그리므로, 이렇게 맞춰야 같은 기준이 된다.
    """
    import re
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from svg2shapes import content_bbox
    x0, y0, x1, y1 = content_bbox(svg)
    return re.sub(r'viewBox\s*=\s*["\'][^"\']+["\']',
                  f'viewBox="{x0:.2f} {y0:.2f} {x1-x0:.2f} {y1-y0:.2f}"',
                  svg, count=1)


def compare(src_svg, group, px_w=900, thr=60):
    """원본 SVG 와 pptx 되돌리기 렌더의 차이. (diff비율, 원본이미지, 되돌린이미지)"""
    from PIL import ImageChops
    back, _ = group_to_svg(group)
    a = render_png(crop_to_content(src_svg), px_w)
    b = render_png(back, px_w)
    if abs(a.size[1] - b.size[1]) > 3:      # 비율이 다르면 그 자체가 오류다
        return 1.0, a, b
    if a.size != b.size:
        b = b.resize(a.size)
    d = ImageChops.difference(a.convert("L"), b.convert("L"))
    n = sum(d.point(lambda v: 255 if v > thr else 0).histogram()[255:])
    return n / float(a.size[0] * a.size[1]), a, b


def check_group_tight(group, tol=1):
    """그룹 경계 == 자식 경계인지. 어긋나면 선택 박스가 그림과 안 맞는다."""
    kids = list(group.shapes)
    if not kids:
        return "빈 그룹"
    x0 = min(s.left for s in kids)
    y0 = min(s.top for s in kids)
    x1 = max(s.left + s.width for s in kids)
    y1 = max(s.top + s.height for s in kids)
    bad = []
    if abs(group.left - x0) > tol or abs(group.top - y0) > tol:
        bad.append(f"off {group.left},{group.top} != {x0},{y0}")
    if abs(group.width - (x1 - x0)) > tol or abs(group.height - (y1 - y0)) > tol:
        bad.append(f"ext {group.width}x{group.height} != {x1-x0}x{y1-y0}")
    xf = group._element.find(qn("p:grpSpPr")).find(qn("a:xfrm"))
    ch, ce = xf.find(qn("a:chOff")), xf.find(qn("a:chExt"))
    if (int(ch.get("x")), int(ch.get("y"))) != (group.left, group.top):
        bad.append("chOff != off (자식이 어긋나 보인다)")
    if (int(ce.get("cx")), int(ce.get("cy"))) != (group.width, group.height):
        bad.append("chExt != ext (자식이 늘어나 보인다)")
    return "; ".join(bad)


if __name__ == "__main__":
    from pptx import Presentation
    from svg2shapes import DEMO, add_svg_shapes
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    g, w, h, n, mn = add_svg_shapes(sl.shapes, DEMO, Emu(914400), Emu(914400),
                                    width=Emu(6 * EMU_PER_IN))
    print("그룹 경계:", check_group_tight(g) or "OK")
    r, a, b = compare(DEMO, g)
    a.save("/tmp/chk_src.png")
    b.save("/tmp/chk_back.png")
    print(f"픽셀 차이 {r*100:.2f}%  (/tmp/chk_src.png, /tmp/chk_back.png)")
